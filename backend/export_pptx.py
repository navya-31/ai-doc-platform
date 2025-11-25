# export_pptx.py
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

def build_pptx_bytes(slides):
    """
    slides: list of dicts {'title':..., 'content':...}
    """
    prs = Presentation()
    # use a simple title + content layout if available
    for s in slides:
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = s.get("title", "")
        # content placeholder is usually placeholder index 1
        try:
            tx = slide.placeholders[1].text = s.get("content", "")
        except Exception:
            # fallback: add a text box
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(4.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.text = s.get("content", "")
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio
